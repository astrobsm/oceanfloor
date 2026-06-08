"""Export Engine.

Renders manuscripts into a variety of file formats. Heavy binary formats
(`docx`, `pptx`, `xlsx`) are produced via lazy imports so the engine still imports
cleanly when those optional packages are missing.

The route handler wraps the returned bytes in a streaming response with the
appropriate content type.
"""
from __future__ import annotations

import csv
import io
import json
from dataclasses import dataclass

from app.schemas.extras import ExportFormat, ExportManuscriptRequest

# Content-type lookup used by the route layer.
MIME_TYPES = {
    ExportFormat.docx: "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ExportFormat.pptx: "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ExportFormat.xlsx: "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ExportFormat.csv: "text/csv",
    ExportFormat.json: "application/json",
    ExportFormat.md: "text/markdown",
    ExportFormat.html: "text/html",
    ExportFormat.latex: "application/x-tex",
}


@dataclass
class ExportArtifact:
    filename: str
    content: bytes
    media_type: str


class ExportEngine:
    def export(self, req: ExportManuscriptRequest) -> ExportArtifact:
        renderer = {
            ExportFormat.docx: self._docx,
            ExportFormat.pptx: self._pptx,
            ExportFormat.xlsx: self._xlsx,
            ExportFormat.csv: self._csv,
            ExportFormat.json: self._json,
            ExportFormat.md: self._md,
            ExportFormat.html: self._html,
            ExportFormat.latex: self._latex,
        }[req.format]
        payload = renderer(req)
        safe_title = "".join(c if c.isalnum() else "_" for c in req.title)[:60] or "manuscript"
        return ExportArtifact(
            filename=f"{safe_title}.{req.format.value}",
            content=payload,
            media_type=MIME_TYPES[req.format],
        )

    # ---- text formats ----
    @staticmethod
    def _md(req: ExportManuscriptRequest) -> bytes:
        parts = [f"# {req.title}\n"]
        for name, body in req.sections.items():
            parts.append(f"## {name}\n\n{body}\n")
        return "\n".join(parts).encode("utf-8")

    @staticmethod
    def _json(req: ExportManuscriptRequest) -> bytes:
        return json.dumps({"title": req.title, "sections": req.sections}, indent=2).encode("utf-8")

    @staticmethod
    def _html(req: ExportManuscriptRequest) -> bytes:
        from html import escape

        body = [f"<h1>{escape(req.title)}</h1>"]
        for name, content in req.sections.items():
            body.append(f"<h2>{escape(name)}</h2><p>{escape(content).replace(chr(10), '<br>')}</p>")
        html = (
            "<!doctype html><html><head><meta charset='utf-8'>"
            f"<title>{escape(req.title)}</title></head><body>"
            + "".join(body)
            + "</body></html>"
        )
        return html.encode("utf-8")

    @staticmethod
    def _latex(req: ExportManuscriptRequest) -> bytes:
        def esc(text: str) -> str:
            for a, b in [("\\", r"\textbackslash{}"), ("&", r"\&"), ("%", r"\%"),
                         ("$", r"\$"), ("#", r"\#"), ("_", r"\_"), ("{", r"\{"), ("}", r"\}")]:
                text = text.replace(a, b)
            return text

        parts = [
            r"\documentclass{article}",
            r"\usepackage[utf8]{inputenc}",
            r"\title{" + esc(req.title) + "}",
            r"\begin{document}", r"\maketitle",
        ]
        for name, body in req.sections.items():
            parts.append(r"\section*{" + esc(name) + "}")
            parts.append(esc(body))
        parts.append(r"\end{document}")
        return "\n".join(parts).encode("utf-8")

    @staticmethod
    def _csv(req: ExportManuscriptRequest) -> bytes:
        buf = io.StringIO()
        writer = csv.writer(buf)
        writer.writerow(["section", "content"])
        for name, body in req.sections.items():
            writer.writerow([name, body])
        return buf.getvalue().encode("utf-8")

    # ---- binary formats (lazy imports keep optional deps optional) ----
    @staticmethod
    def _docx(req: ExportManuscriptRequest) -> bytes:
        try:
            from docx import Document  # type: ignore
        except ImportError as exc:
            raise RuntimeError("python-docx is not installed") from exc
        doc = Document()
        doc.add_heading(req.title, level=0)
        for name, body in req.sections.items():
            doc.add_heading(name, level=1)
            doc.add_paragraph(body)
        out = io.BytesIO()
        doc.save(out)
        return out.getvalue()

    @staticmethod
    def _pptx(req: ExportManuscriptRequest) -> bytes:
        try:
            from pptx import Presentation  # type: ignore
            from pptx.util import Inches
        except ImportError as exc:
            raise RuntimeError("python-pptx is not installed") from exc

        prs = Presentation()
        title_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_layout)
        slide.shapes.title.text = req.title

        content_layout = prs.slide_layouts[1]
        for name, body in req.sections.items():
            s = prs.slides.add_slide(content_layout)
            s.shapes.title.text = name
            tf = s.placeholders[1].text_frame
            for line in body.splitlines() or [""]:
                tf.add_paragraph().text = line[:200]
        out = io.BytesIO()
        prs.save(out)
        return out.getvalue()

    @staticmethod
    def _xlsx(req: ExportManuscriptRequest) -> bytes:
        try:
            from openpyxl import Workbook  # type: ignore
        except ImportError as exc:
            raise RuntimeError("openpyxl is not installed") from exc
        wb = Workbook()
        ws = wb.active
        ws.title = "Manuscript"
        ws.append(["Section", "Content"])
        for name, body in req.sections.items():
            ws.append([name, body])
        out = io.BytesIO()
        wb.save(out)
        return out.getvalue()


export_engine = ExportEngine()
